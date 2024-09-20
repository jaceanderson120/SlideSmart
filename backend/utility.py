from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from flask import Flask, jsonify
from flask_cors import CORS
import io
import base64
from gpt import analyze_powerpoint

app = Flask(__name__)
CORS(app)

def extract_text_and_images_from_ppt(ppt):
    """Extract all text and images from a PowerPoint and return a dictionary.

    Arguments:
    - ppt: the PowerPoint to parse

    Returns:
    - a dictionary of the following form: {"slide1": {"text": [text1, text2, ...], "images": [image1, image2, ...]}, "slide2": {...}, ...}
    """
    
    slides_data = {}

    for slide_number, slide in enumerate(ppt.slides, start=1):
        text_array = []
        image_array = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Convert images to base64 so that they can be jsonify'd
                image_stream = io.BytesIO(shape.image.blob)
                base64_image = base64.b64encode(image_stream.read()).decode('utf-8')
                image_array.append(base64_image)
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_array.append(run.text)

        # Combine text and images into a single dictionary for each slide
        slide_key = f'slide{slide_number}'
        slides_data[slide_key] = {
            "text": text_array,
            "images": [] # Appending an empty image array for now
        }

    return analyze_powerpoint(slides_data)

@app.route('/extract', methods=['POST'])
def upload_file():
    file_path = "01-overview.pptx" # Currently this is a PowerPoint in my local root directory, change this to accept PowerPoints from the frontend
    ppt = Presentation(file_path)
    extracted_data = extract_text_and_images_from_ppt(ppt)
    return jsonify(extracted_data)

if __name__ == "__main__":
    app.run(debug=True)
